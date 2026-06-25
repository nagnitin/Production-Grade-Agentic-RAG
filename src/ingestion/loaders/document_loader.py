"""
Unified document loader factory.

WHY: Different file formats require different parsing strategies. This factory
provides a single interface that dispatches to the appropriate LangChain loader
based on file extension. All loaders return a common List[Document] format.

ARCHITECTURE DECISION: Using LangChain's built-in loaders instead of building
custom parsers because:
1. Battle-tested across thousands of production deployments
2. Consistent Document output format with metadata
3. Active maintenance and bug fixes
4. Easy to swap or extend individual loaders

TRADE-OFF: LangChain loaders require files on disk (not in-memory bytes).
We write temp files, which adds ~5ms overhead per document. Acceptable for
a batch ingestion pipeline where parsing takes 100ms–10s per document.
"""

from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import Any

from langchain_core.documents import Document

from src.config.logging_config import get_logger

logger = get_logger(__name__)


class DocumentLoader:
    """
    Factory for loading documents from various file formats.

    Supports: PDF, DOCX, PPTX, HTML, TXT
    All loaders return List[Document] with standardized metadata.
    """

    # Supported extensions and their loader factories
    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".html", ".htm", ".txt"}

    def load_from_bytes(
        self,
        content: bytes,
        filename: str,
        file_type: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[Document]:
        """
        Load documents from raw bytes.

        Writes to a temp file because most LangChain loaders require
        file paths. The temp file is cleaned up after loading.

        Args:
            content: Raw file bytes
            filename: Original filename
            file_type: File extension (e.g., ".pdf")
            metadata: Additional metadata to attach to all documents

        Returns:
            List of LangChain Documents with page content and metadata
        """
        metadata = metadata or {}
        file_type = file_type.lower()

        if file_type not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: {file_type}. "
                f"Supported: {self.SUPPORTED_EXTENSIONS}"
            )

        logger.info(
            "Loading document",
            filename=filename,
            file_type=file_type,
            size_bytes=len(content),
        )

        # Write to temp file for LangChain loaders
        with tempfile.NamedTemporaryFile(
            suffix=file_type, delete=False
        ) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            documents = self._dispatch_loader(tmp_path, file_type)
        finally:
            # Clean up temp file
            os.unlink(tmp_path)

        # Enrich metadata on all documents
        base_metadata = {
            "source": filename,
            "file_type": file_type,
            "original_size_bytes": len(content),
            **metadata,
        }

        for doc in documents:
            doc.metadata = {**base_metadata, **doc.metadata}
            doc.metadata["source"] = filename

        logger.info(
            "Document loaded",
            filename=filename,
            num_pages=len(documents),
        )

        return documents

    def load_from_path(
        self,
        file_path: str | Path,
        metadata: dict[str, Any] | None = None,
    ) -> list[Document]:
        """Load documents from a file path on disk."""
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_type = path.suffix.lower()
        content = path.read_bytes()

        return self.load_from_bytes(
            content=content,
            filename=path.name,
            file_type=file_type,
            metadata=metadata,
        )

    def _dispatch_loader(
        self, file_path: str, file_type: str
    ) -> list[Document]:
        """Dispatch to the appropriate loader based on file type."""
        loaders = {
            ".pdf": self._load_pdf,
            ".docx": self._load_docx,
            ".pptx": self._load_pptx,
            ".html": self._load_html,
            ".htm": self._load_html,
            ".txt": self._load_txt,
        }

        loader_fn = loaders.get(file_type)
        if not loader_fn:
            raise ValueError(f"No loader for file type: {file_type}")

        return loader_fn(file_path)

    @staticmethod
    def _load_pdf(file_path: str) -> list[Document]:
        """
        Load PDF using PyPDF.

        Returns one Document per page with page number metadata.
        Falls back to unstructured loader on failure.
        """
        try:
            from langchain_community.document_loaders import PyPDFLoader

            loader = PyPDFLoader(file_path, extract_images=False)
            documents = loader.load()

            # Enrich with page numbers (0-indexed → 1-indexed)
            for doc in documents:
                if "page" in doc.metadata:
                    doc.metadata["page"] = doc.metadata["page"] + 1

            return documents

        except Exception as e:
            logger.warning(
                "PyPDF failed, trying unstructured fallback",
                error=str(e),
            )
            return DocumentLoader._load_unstructured(file_path)

    @staticmethod
    def _load_docx(file_path: str) -> list[Document]:
        """Load DOCX using Docx2txt."""
        try:
            from langchain_community.document_loaders import Docx2txtLoader

            loader = Docx2txtLoader(file_path)
            return loader.load()

        except ImportError:
            logger.warning("Docx2txt not available, using unstructured")
            return DocumentLoader._load_unstructured(file_path)

    @staticmethod
    def _load_pptx(file_path: str) -> list[Document]:
        """
        Load PPTX using python-pptx.

        Extracts text from each slide as a separate document,
        preserving slide number metadata.
        """
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            documents = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                text_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                text_parts.append(text)

                if text_parts:
                    content = "\n".join(text_parts)
                    documents.append(
                        Document(
                            page_content=content,
                            metadata={
                                "slide_number": slide_num,
                                "page": slide_num,
                            },
                        )
                    )

            return documents

        except ImportError:
            logger.warning("python-pptx not available, using unstructured")
            return DocumentLoader._load_unstructured(file_path)

    @staticmethod
    def _load_html(file_path: str) -> list[Document]:
        """Load HTML using BeautifulSoup."""
        try:
            from langchain_community.document_loaders import BSHTMLLoader

            loader = BSHTMLLoader(
                file_path,
                open_encoding="utf-8",
                get_text_separator="\n",
            )
            return loader.load()

        except Exception as e:
            logger.warning("BSHTMLLoader failed", error=str(e))
            # Fallback: read as plain text
            return DocumentLoader._load_txt(file_path)

    @staticmethod
    def _load_txt(file_path: str) -> list[Document]:
        """
        Load plain text files.

        Handles encoding detection for non-UTF-8 files.
        """
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
        except UnicodeDecodeError:
            # Try detecting encoding
            import chardet

            with open(file_path, "rb") as f:
                raw = f.read()

            detected = chardet.detect(raw)
            encoding = detected.get("encoding", "latin-1")
            content = raw.decode(encoding, errors="replace")

        return [Document(page_content=content, metadata={})]

    @staticmethod
    def _load_unstructured(file_path: str) -> list[Document]:
        """
        Fallback loader using the unstructured library.

        Handles almost any document format but is slower and
        requires more dependencies.
        """
        try:
            from langchain_community.document_loaders import (
                UnstructuredFileLoader,
            )

            loader = UnstructuredFileLoader(
                file_path,
                mode="elements",
                strategy="fast",
            )
            return loader.load()

        except Exception as e:
            logger.error(
                "All loaders failed for document",
                error=str(e),
                file_path=file_path,
            )
            raise RuntimeError(
                f"Failed to load document: {e}"
            ) from e
